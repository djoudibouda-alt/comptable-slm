#!/usr/bin/env python3
"""
Génère le PPTX Comptable_SLM_Pitch.pptx à partir des slides PNG.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT_PPTX = Path("Comptable_SLM_Pitch.pptx")
SLIDES_DIR = Path("slides")

# Couleurs
DARK_BG = RGBColor(0x0f, 0x17, 0x2a)   # slate-900
ACCENT = RGBColor(0x38, 0xbd, 0xf8)    # sky-400
ACCENT2 = RGBColor(0x22, 0xd3, 0xee)   # cyan-400
WHITE = RGBColor(0xf8, 0xfa, 0xfc)
GRAY = RGBColor(0x94, 0xa3, 0xb8)
GREEN = RGBColor(0x22, 0xc5, 0x5e)

prs = Presentation()
prs.slide_width = Inches(10.8)
prs.slide_height = Inches(13.5)

# Layout vide
blank_layout = prs.slide_layouts[6]  # Blank

slide_files = [
    ("01_couverture.png", "Couverture"),
    ("02_probleme.png", "Le problème"),
    ("03_solution.png", "La solution"),
    ("04_architecture.png", "Architecture"),
    ("05_knowledgebase.png", "Base de connaissances"),
    ("06_qa_curated.png", "125 Q&A curées"),
    ("07_resultats.png", "Résultats tests"),
    ("08_demo.png", "Démo en ligne"),
    ("09_roadmap.png", "Roadmap & Vision"),
    ("10_cta.png", "Call to Action"),
]

for i, (fname, title) in enumerate(slide_files):
    slide = prs.slides.add_slide(blank_layout)
    
    # Fond
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    # Image
    img_path = SLIDES_DIR / fname
    if img_path.exists():
        slide.shapes.add_picture(
            str(img_path),
            Inches(0), Inches(0),
            width=Inches(10.8), height=Inches(13.5)
        )
    else:
        # Fallback texte si image manquante
        txBox = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8.8), Inches(3))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"Slide {i+1}: {title}"
        p.font.size = Pt(36)
        p.font.color.rgb = ACCENT
        p.alignment = PP_ALIGN.CENTER

prs.save(OUT_PPTX)
print(f"[OK] PPTX genere : {OUT_PPTX} ({len(slide_files)} slides)")